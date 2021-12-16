# Make it run from the examples directory
import sys
sys.path.append("..")
from liquer import *
from liquer.query import evaluate, evaluate_and_save
from liquer.store import set_store, MemoryStore
import importlib
from pptx import Presentation

import liquer.ext.lq_pptx  # register pptx commands and state type
import liquer.ext.basic  # dr command


@first_command
def make_presentation():
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "Hello, World!"
    subtitle.text = "python-pptx was here!"
    return prs

@command
def add_slide(prs,title="Title"):
    bullet_slide_layout = prs.slide_layouts[1]

    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes

    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = title

    tf = body_shape.text_frame
    tf.text = 'Find the bullet slide layout'

    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.text for first bullet'
    p.level = 1
    return prs


evaluate_and_save(f"make_presentation/add_slide-Slide1/add_slide-Slide2/test.pptx")
 
